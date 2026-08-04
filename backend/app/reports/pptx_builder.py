"""PowerPoint report generation using python-pptx"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime
import logging

logger = logging.getLogger(__name__)


def generate_pptx_report(
    question: str,
    executive_summary: dict,
    kpis: dict,
    recommendations: list,
    forecasts: dict = None,
    output_path: str = "/tmp/report.pptx",
) -> str:
    """
    Generate a PowerPoint report with multiple slides.
    """
    try:
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Color scheme
        title_color = RGBColor(31, 78, 120)  # #1f4e78
        accent_color = RGBColor(46, 92, 138)  # #2e5c8a

        # Slide 1: Title Slide
        slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        slide1.background.fill.solid()
        slide1.background.fill.fore_color.rgb = RGBColor(240, 240, 240)

        title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = "DecisionLens AI"
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = title_color
        p.alignment = PP_ALIGN.CENTER

        subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1.5))
        subtitle_frame = subtitle_box.text_frame
        p = subtitle_frame.paragraphs[0]
        p.text = f"Business Analysis Report\n{datetime.utcnow().strftime('%B %d, %Y')}"
        p.font.size = Pt(18)
        p.font.color.rgb = accent_color
        p.alignment = PP_ALIGN.CENTER

        # Slide 2: Question & Executive Summary
        slide2 = prs.slides.add_slide(prs.slide_layouts[6])
        slide2.background.fill.solid()
        slide2.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

        # Question
        q_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        q_frame = q_box.text_frame
        q_frame.word_wrap = True
        p = q_frame.paragraphs[0]
        p.text = f"Question: {question}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = title_color

        # Executive Summary
        if executive_summary and isinstance(executive_summary, dict):
            summary_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(9), Inches(5.5))
            summary_frame = summary_box.text_frame
            summary_frame.word_wrap = True

            if "narrative" in executive_summary:
                p = summary_frame.paragraphs[0]
                p.text = executive_summary["narrative"][:500]  # First 500 chars
                p.font.size = Pt(12)

        # Slide 3: KPIs
        if kpis:
            slide3 = prs.slides.add_slide(prs.slide_layouts[6])
            slide3.background.fill.solid()
            slide3.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

            # Title
            title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
            title_frame = title_box.text_frame
            p = title_frame.paragraphs[0]
            p.text = "Key Performance Indicators"
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = title_color

            # KPI boxes
            kpi_items = list(kpis.items())[:6]
            for i, (metric, info) in enumerate(kpi_items):
                row = i // 3
                col = i % 3

                x = 0.5 + col * 3.2
                y = 1.2 + row * 2.8

                box = slide3.shapes.add_shape(
                    1,  # Rectangle shape
                    Inches(x),
                    Inches(y),
                    Inches(3),
                    Inches(2.3)
                )
                box.fill.solid()
                box.fill.fore_color.rgb = RGBColor(230, 240, 250)
                box.line.color.rgb = accent_color

                text_frame = box.text_frame
                text_frame.word_wrap = True

                p = text_frame.paragraphs[0]
                p.text = metric
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = title_color

                if isinstance(info, dict):
                    value_text = f"{info.get('value', 'N/A')}"
                    if isinstance(info.get('value'), (int, float)):
                        value_text = f"{float(info['value']):.2f}"

                    p = text_frame.add_paragraph()
                    p.text = value_text
                    p.font.size = Pt(16)
                    p.font.bold = True
                    p.font.color.rgb = accent_color
                    p.space_before = Pt(4)

        # Slide 4: Recommendations
        if recommendations:
            slide4 = prs.slides.add_slide(prs.slide_layouts[6])
            slide4.background.fill.solid()
            slide4.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

            title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
            title_frame = title_box.text_frame
            p = title_frame.paragraphs[0]
            p.text = "Recommended Actions"
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = title_color

            rec_box = slide4.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
            rec_frame = rec_box.text_frame
            rec_frame.word_wrap = True

            for i, rec in enumerate(recommendations[:5], 1):
                title_text = rec.get("title", f"Action {i}") if isinstance(rec, dict) else f"Action {i}"
                priority = rec.get("priority", "") if isinstance(rec, dict) else ""

                p = rec_frame.add_paragraph() if i > 1 else rec_frame.paragraphs[0]
                p.text = f"{i}. {title_text} [{priority}]"
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = accent_color
                p.space_before = Pt(6)

                if isinstance(rec, dict) and "expected_impact" in rec:
                    p = rec_frame.add_paragraph()
                    p.text = f"   Impact: {rec['expected_impact']}"
                    p.font.size = Pt(10)
                    p.level = 1

        # Save presentation
        prs.save(output_path)

        logger.info(f"PowerPoint report generated: {output_path}")
        return output_path

    except Exception as e:
        logger.error(f"PowerPoint generation failed: {e}")
        raise
